"""
Extract plain text from supported file types and perform in-place token
substitution for office formats during obfuscation / rehydration.
"""
from __future__ import annotations

import io
import re as _re
from typing import BinaryIO, Callable

OFFICE_EXTS = {".docx", ".pptx", ".xlsx"}


# ── extraction ──────────────────────────────────────────────────────────────

def extract_text(file_obj: BinaryIO, ext: str) -> str:
    if ext == ".docx":
        return _extract_docx(file_obj)
    if ext == ".pptx":
        return _extract_pptx(file_obj)
    if ext == ".xlsx":
        return _extract_xlsx(file_obj)
    raw = file_obj.read()
    return raw.decode("utf-8", errors="replace")


def _extract_docx(file_obj: BinaryIO) -> str:
    from docx import Document
    doc = Document(file_obj)
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    return "\n".join(parts)


def _extract_pptx(file_obj: BinaryIO) -> str:
    from pptx import Presentation
    prs = Presentation(file_obj)
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text_frame.text.strip()
                        if t:
                            parts.append(t)
    return "\n".join(parts)


def _extract_xlsx(file_obj: BinaryIO) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(file_obj, read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    parts.append(str(cell.value))
    wb.close()
    return "\n".join(parts)


# ── obfuscation (in-place term→token substitution) ─────────────────────────

def obfuscate_file(file_obj: BinaryIO, ext: str, dictionary: list[dict]) -> bytes:
    sorted_entries = sorted(dictionary, key=lambda e: len(e["original_term"]), reverse=True)

    def replace(text: str) -> str:
        for entry in sorted_entries:
            pat = _re.compile(_re.escape(entry["original_term"]), _re.IGNORECASE)
            text = pat.sub(entry["token"], text)
        return text

    if ext == ".docx":
        return _docx_apply(file_obj, replace)
    if ext == ".pptx":
        return _pptx_apply(file_obj, replace)
    if ext == ".xlsx":
        return _xlsx_apply(file_obj, replace)
    raise ValueError(f"Unsupported office format for obfuscation: {ext}")


# ── rehydration (in-place token→term substitution) ─────────────────────────

def rehydrate_file(file_obj: BinaryIO, ext: str, dictionary: list[dict]) -> bytes:
    def replace(text: str) -> str:
        for entry in dictionary:
            text = text.replace(entry["token"], entry["original_term"])
        return text

    if ext == ".docx":
        return _docx_apply(file_obj, replace)
    if ext == ".pptx":
        return _pptx_apply(file_obj, replace)
    if ext == ".xlsx":
        return _xlsx_apply(file_obj, replace)
    raise ValueError(f"Unsupported office format for rehydration: {ext}")


# ── format-specific helpers (shared between obfuscate / rehydrate) ─────────

def _apply_runs(para, replace_fn: Callable[[str], str]) -> None:
    if not para.runs:
        return
    full = para.text
    replaced = replace_fn(full)
    if replaced != full:
        para.runs[0].text = replaced
        for run in para.runs[1:]:
            run.text = ""


def _docx_apply(file_obj: BinaryIO, replace_fn: Callable[[str], str]) -> bytes:
    from docx import Document
    doc = Document(file_obj)
    for para in doc.paragraphs:
        _apply_runs(para, replace_fn)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _apply_runs(para, replace_fn)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_apply(file_obj: BinaryIO, replace_fn: Callable[[str], str]) -> bytes:
    from pptx import Presentation
    prs = Presentation(file_obj)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    _apply_runs(para, replace_fn)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            _apply_runs(para, replace_fn)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _xlsx_apply(file_obj: BinaryIO, replace_fn: Callable[[str], str]) -> bytes:
    from openpyxl import load_workbook
    wb = load_workbook(file_obj)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    cell.value = replace_fn(cell.value)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
